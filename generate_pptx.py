from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_FILE = r"e:\Project\jansewa-ai\Jansewa_AI_Architecture_Evaluation.pptx"

# ---------- Theme ----------
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(20, 20, 20)
DARK_BLUE = RGBColor(22, 56, 92)
MID_GRAY = RGBColor(90, 90, 90)
LIGHT_GRAY = RGBColor(235, 238, 242)

TITLE_SIZE = Pt(34)
SUBTITLE_SIZE = Pt(18)
HEADING_SIZE = Pt(28)
BODY_SIZE = Pt(18)
SMALL_SIZE = Pt(14)

prs = Presentation()  # default 16:9


def set_slide_bg_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_top_bar(slide, title_text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.size = Pt(18)
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.3)


def add_footer(slide, text="Jansewa AI • Sankalp Hackathon 2025–2026"):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.35), prs.slide_width, Inches(0.35))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_GRAY
    line.line.fill.background()

    tf = line.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = SMALL_SIZE
    run.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.CENTER


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg_white(slide)

    # Title block
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.4))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = TITLE_SIZE
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE

    sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.2), Inches(11.0), Inches(1.2))
    tf2 = sub_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = SUBTITLE_SIZE
    run2.font.color.rgb = MID_GRAY

    # Accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(4.2), Inches(5.5), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = DARK_BLUE
    accent.line.fill.background()

    add_footer(slide)
    return slide


def add_bullets_slide(slide_title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_white(slide)
    add_top_bar(slide, slide_title)

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.1), Inches(11.5), Inches(5.8))
    tf = body.text_frame
    tf.clear()

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = BODY_SIZE
        p.font.color.rgb = BLACK

    add_footer(slide)
    return slide


def add_two_column_slide(slide_title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_white(slide)
    add_top_bar(slide, slide_title)

    # Left panel
    lpanel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.1), Inches(5.8), Inches(5.8))
    lpanel.fill.solid()
    lpanel.fill.fore_color.rgb = RGBColor(248, 250, 253)
    lpanel.line.color.rgb = LIGHT_GRAY

    ltitle = slide.shapes.add_textbox(Inches(1.0), Inches(1.35), Inches(5.2), Inches(0.5))
    ltf = ltitle.text_frame
    ltf.text = left_title
    ltf.paragraphs[0].font.size = Pt(22)
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.color.rgb = DARK_BLUE

    lbody = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(5.2), Inches(4.6))
    lbtf = lbody.text_frame
    lbtf.clear()
    for i, item in enumerate(left_items):
        p = lbtf.paragraphs[0] if i == 0 else lbtf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK

    # Right panel
    rpanel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.1), Inches(5.8), Inches(5.8))
    rpanel.fill.solid()
    rpanel.fill.fore_color.rgb = RGBColor(248, 250, 253)
    rpanel.line.color.rgb = LIGHT_GRAY

    rtitle = slide.shapes.add_textbox(Inches(7.1), Inches(1.35), Inches(5.2), Inches(0.5))
    rtf = rtitle.text_frame
    rtf.text = right_title
    rtf.paragraphs[0].font.size = Pt(22)
    rtf.paragraphs[0].font.bold = True
    rtf.paragraphs[0].font.color.rgb = DARK_BLUE

    rbody = slide.shapes.add_textbox(Inches(7.1), Inches(1.9), Inches(5.2), Inches(4.6))
    rbtf = rbody.text_frame
    rbtf.clear()
    for i, item in enumerate(right_items):
        p = rbtf.paragraphs[0] if i == 0 else rbtf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK

    add_footer(slide)
    return slide


def add_architecture_flow_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_white(slide)
    add_top_bar(slide, "Architecture Flow (Easy View)")

    # Boxes
    y = Inches(2.2)
    box_w = Inches(2.25)
    box_h = Inches(1.2)
    x_positions = [Inches(0.6), Inches(3.2), Inches(5.8), Inches(8.4), Inches(11.0)]
    labels = [
        "Users\n(Citizen/Leader)",
        "Frontend\n(React)",
        "Backend API\n(FastAPI)",
        "Services + KB\n(Offline AI)",
        "DB/Cache\n(Postgres/Redis)",
    ]

    for x, label in zip(x_positions, labels):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(248, 250, 253)
        shp.line.color.rgb = DARK_BLUE
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLACK

    # Arrows
    for i in range(4):
        ax = x_positions[i] + box_w
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax + Inches(0.05), y + Inches(0.35), Inches(0.35), Inches(0.45))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_BLUE
        arrow.line.fill.background()

    caption = slide.shapes.add_textbox(Inches(0.9), Inches(4.2), Inches(11.6), Inches(2.2))
    tfc = caption.text_frame
    tfc.text = "KB-first logic: Local Knowledge Base runs first. Optional APIs (Gemini/Whisper/Twitter) only enhance output."
    tfc.paragraphs[0].font.size = Pt(20)
    tfc.paragraphs[0].font.color.rgb = DARK_BLUE
    tfc.paragraphs[0].font.bold = True

    p2 = tfc.add_paragraph()
    p2.text = "This ensures reliability in low-connectivity environments and zero API-key dependency for core features."
    p2.font.size = Pt(16)
    p2.font.color.rgb = BLACK

    add_footer(slide)


# ---------- Build Slides ----------
add_title_slide(
    "Jansewa AI",
    "AI-Powered Real-Time Governance Intelligence Platform\nArchitecture & Evaluation Presentation"
)

add_bullets_slide("Problem & Vision", [
    "Local governance teams need fast, transparent, and data-driven complaint management.",
    "Jansewa AI converts citizen complaints into prioritized, verifiable civic actions.",
    "Vision: Reliable AI for public governance, even with limited internet/API access."
])

add_two_column_slide(
    "System Overview",
    "Frontend (User Interaction)",
    [
        "React + Vite + Tailwind dashboard",
        "Complaint intake, verification view, social insights",
        "Heatmap, charts, trust score, public portal"
    ],
    "Backend (Execution Engine)",
    [
        "FastAPI with domain-based routers",
        "JWT auth with role-based access",
        "AI services + Knowledge Base integration"
    ]
)

add_architecture_flow_slide()

add_bullets_slide("Knowledge Base (Core Innovation)", [
    "8 self-contained modules provide offline intelligence.",
    "Covers complaint classification, priority rules, SOPs, templates, ward profiles.",
    "Works without external APIs; optional AI only improves output quality.",
    "Design principle: KB PRIMARY, external APIs OPTIONAL."
])

add_bullets_slide("Complaint Processing Pipeline", [
    "1) Citizen submits text/voice/image complaint.",
    "2) KB classifier identifies category, urgency, and possible duplicates.",
    "3) Priority engine scores using 5 factors + escalation rules.",
    "4) Department routing and task assignment are generated.",
    "5) Complaint lifecycle tracked with status and audit logs."
])

add_bullets_slide("4-Layer Verification Architecture", [
    "Layer 1: GPS consistency check",
    "Layer 2: Timestamp validation",
    "Layer 3: Image comparison (pixel-diff / vision enhancement optional)",
    "Layer 4: Tamper detection and integrity checks",
    "Leader approves/rejects for accountability and closure quality."
])

add_two_column_slide(
    "Data & Analytics Layer",
    "Transactional Data (PostgreSQL)",
    [
        "Complaints, users, wards, verifications, communications",
        "Trust scores and social analysis records",
        "Audit logs for accountability"
    ],
    "Operational Intelligence",
    [
        "Priority queue for leadership decisions",
        "Ward heatmap and sentiment trends",
        "Public scorecards for transparency"
    ]
)

add_bullets_slide("Security, Access & Governance", [
    "JWT-based authentication with role-specific permissions.",
    "Roles: Leader, Department Head, Worker, Admin.",
    "Approval workflows for verification and communications.",
    "Audit trails for traceability and governance compliance."
])

add_bullets_slide("Deployment & DevOps", [
    "Docker Compose runs full stack: frontend, backend, postgres, redis.",
    "CI/CD with GitHub Actions: lint, test, build, deploy.",
    "Backend deploy target: Railway; Frontend: Vercel; DB: Supabase/Postgres.",
    "Fast evaluation via seeded demo data and Swagger docs."
])

add_bullets_slide("Why This Architecture Is Evaluation-Ready", [
    "Offline-first reliability for public sector constraints.",
    "Modular and maintainable service-based backend design.",
    "Clear accountability with verification + audit model.",
    "Citizen trust focus through transparency portal and trust metrics.",
    "Scalable foundation for multi-ward and multi-city rollout."
])

add_bullets_slide("Demo Flow (5 Minutes)", [
    "1) Create complaint (text/voice).",
    "2) Show auto classification + priority score.",
    "3) Display dashboard heatmap and queue.",
    "4) Submit and approve verification evidence.",
    "5) Generate and publish bilingual communication.",
    "6) Open public portal scorecard for transparency."
])

# Save
prs.save(OUTPUT_FILE)
print(f"Presentation generated: {OUTPUT_FILE}")